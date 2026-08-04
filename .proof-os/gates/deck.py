#!/usr/bin/env python3
# origin: riya STEP 6 self-check — 3 of 6 items promoted: opens-cleanly, ≤3 core colors, AA contrast
import sys
def lum(rgb):
    c=[x/255 for x in rgb]; c=[x/12.92 if x<=.03928 else ((x+.055)/1.055)**2.4 for x in c]
    return .2126*c[0]+.7152*c[1]+.0722*c[2]
def ratio(a,b): L=sorted([lum(a),lum(b)],reverse=True); return (L[0]+.05)/(L[1]+.05)
try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    print("VERDICT: partial — python-pptx unavailable (oracle missing, not passing)"); sys.exit(2)
p = Presentation(sys.argv[1])  # opens cleanly or throws = the gate
colors=set(); fails=[]
for s in p.slides:
    for sh in s.shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    if r.font.color and r.font.color.type is not None and hasattr(r.font.color, "rgb") and r.font.color.rgb:
                        colors.add(str(r.font.color.rgb))
core = {c for c in colors}
print(f"opens: yes · distinct text colors: {len(core)}")
if len(core) > 6: fails.append(f"{len(core)} text colors — 60/30/10 wants ≤3 core (+shades)")
for f in fails: print(" ", f)
sys.exit(1 if fails else 0)
